"""Update PPT with sprint code snapshot slides."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_IN = (
    "d:/Dhanaraj/Learning/Agentic AI/NIIT/Build and Deploy Intelligent"
    " Conversational Agents/Sprint 11/course2-capstoneproject/docs"
    "/Job_Placement_Agent_Demo.pptx"
)
PPTX = (
    "d:/Dhanaraj/Learning/Agentic AI/NIIT/Build and Deploy Intelligent"
    " Conversational Agents/Sprint 11/course2-capstoneproject/docs"
    "/Job_Placement_Agent_Demo_v2.pptx"
)

# ── Colours ───────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
C_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)
C_COMMENT = RGBColor(0x6A, 0x99, 0x55)
C_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)
C_DECO    = RGBColor(0xDC, 0xDC, 0xAA)
C_STRING  = RGBColor(0xCE, 0x91, 0x78)
C_HEADER  = RGBColor(0x8A, 0xC8, 0xFF)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

KWDS = ('def ', 'async def ', 'class ', 'return ', 'raise ',
        'from ', 'import ', 'try:', 'except ', 'if ', 'else:',
        'elif ', 'for ', 'with ', 'await ', 'yield ', 'lambda ')

def line_color(line):
    s = line.strip()
    if not s:                         return C_DEFAULT
    if s.startswith('#'):             return C_COMMENT
    if s.startswith('@'):             return C_DECO
    if '"""' in s or "'''" in s:     return C_STRING
    if any(s.startswith(k) for k in KWDS): return C_KEYWORD
    return C_DEFAULT


def add_rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_tb(slide, x, y, w, h, text, size, color,
           bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return tb


def add_code(slide, x, y, w, h, lines, size=Pt(9.5)):
    add_rect(slide, x, y, w, h, C_CODE_BG)
    tb = slide.shapes.add_textbox(
        x + Inches(0.08), y + Inches(0.06),
        w - Inches(0.16), h - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln if ln else ' '
        r.font.size = size
        r.font.name = "Courier New"
        r.font.color.rgb = line_color(ln)


def insert_blank_slide(prs, after_index):
    layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    lst = prs.slides._sldIdLst
    el = lst[-1]
    lst.remove(el)
    lst.insert(after_index + 1, el)
    return prs.slides[after_index + 1]


# ── Sprint data ───────────────────────────────────────────────────────────────
SPRINTS = [
    dict(
        num="Sprint 2", topic="LLM Integration",
        subtitle="ChatOpenAI via Gemini OpenAI-compat API  \u00b7  lru_cache singleton  \u00b7  Async .ainvoke()",
        file="src/backend/agent/llm.py",
        accent=RGBColor(0x23, 0x6B, 0xD4),
        code=[
            "# agent/llm.py \u2014 Sprint 2: LLM Integration via OpenAI-compatible Gemini API",
            "from langchain_openai import ChatOpenAI",
            "from functools import lru_cache",
            "",
            "@lru_cache(maxsize=1)            # singleton \u2014 built once, reused forever",
            "def get_llm():",
            "    api_key  = os.getenv('GEMINI_API_KEY', '')",
            "    model    = os.getenv('GEMINI_MODEL_NAME', 'gemini-2.0-flash')",
            "    base_url = os.getenv('GEMINI_BASE_URL',",
            "        'https://generativelanguage.googleapis.com/v1beta/openai/')",
            "    if not api_key:",
            "        raise GeminiConfigError('GEMINI_API_KEY is not set.')",
            "    return ChatOpenAI(model=model, api_key=api_key,",
            "                     base_url=base_url, temperature=0.7, max_tokens=4096)",
            "",
            "# Sprint 2, LO3 \u2014 Asynchronous LLM call with .ainvoke()",
            "async def run_agent_async(user_message, session_id='default'):",
            "    result = await _get_agent().ainvoke(   # non-blocking async call",
            "        {'input': user_message},",
            "        config={'configurable': {'session_id': session_id}},",
            "    )",
            "    return result.get('output', '')",
        ],
    ),
    dict(
        num="Sprint 3", topic="LangChain Tools",
        subtitle="search_jobs  \u00b7  analyze_resume  \u00b7  generate_cover_letter  \u2014  @tool decorator",
        file="src/backend/agent/tools/job_search.py  |  resume_analyzer.py  |  cover_letter.py",
        accent=RGBColor(0x2E, 0xA0, 0x43),
        code=[
            "# agent/tools/job_search.py \u2014 Sprint 3: LangChain @tool decorator",
            "from langchain.tools import tool",
            "",
            "@tool",
            "def search_jobs(query: str, location: str = '') -> str:",
            '    """Search for live job listings via SerpAPI Google Jobs.',
            "    Use this when the user wants to find job opportunities.",
            '    """',
            "    try:",
            "        return fetch_job_listings(query, location)",
            "    except SerpApiConfigError as exc:",
            "        return f'Configuration error: {exc}'",
            "    except SerpApiRateLimitError:",
            "        return 'Rate limit hit. Please try again in a few minutes.'",
            "    except SerpApiNetworkError:",
            "        return 'Cannot reach SerpAPI. Check your connection.'",
            "    except SerpApiError as exc:",
            "        return f'Job search error: {exc}'",
            "",
            "# All three @tool functions registered with AgentExecutor:",
            "tools = [search_jobs, analyze_resume, generate_cover_letter]",
        ],
    ),
    dict(
        num="Sprint 4", topic="Agent + Memory",
        subtitle="create_tool_calling_agent  \u00b7  AgentExecutor  \u00b7  RunnableWithMessageHistory",
        file="src/backend/agent/job_agent.py",
        accent=RGBColor(0x9B, 0x59, 0xB6),
        code=[
            "# agent/job_agent.py \u2014 Sprint 4: Tool-calling agent with session memory",
            "from langchain.agents import AgentExecutor, create_tool_calling_agent",
            "from langchain_core.runnables.history import RunnableWithMessageHistory",
            "",
            "def _build_agent_executor() -> AgentExecutor:",
            "    prompt = ChatPromptTemplate.from_messages([",
            "        ('system', SYSTEM_PROMPT),",
            "        MessagesPlaceholder('chat_history'),  # injected by history wrapper",
            "        ('human', '{input}'),",
            "        MessagesPlaceholder('agent_scratchpad'),",
            "    ])",
            "    agent    = create_tool_calling_agent(get_llm(), tools, prompt)",
            "    return AgentExecutor(agent=agent, tools=tools,",
            "                        max_iterations=6, handle_parsing_errors=True)",
            "",
            "# Per-session conversation memory (in-process dict)",
            "_session_store: dict[str, ChatMessageHistory] = {}",
            "",
            "agent_with_history = RunnableWithMessageHistory(",
            "    executor, get_session_history,",
            "    input_messages_key='input', history_messages_key='chat_history',",
            ")",
        ],
    ),
    dict(
        num="Sprint 5", topic="LCEL Chains",
        subtitle="Pipe syntax  prompt | llm | StrOutputParser()  in every tool",
        file="src/backend/agent/tools/resume_analyzer.py",
        accent=RGBColor(0xE6, 0x7E, 0x22),
        code=[
            "# agent/tools/resume_analyzer.py \u2014 Sprint 5: LCEL chain with pipe syntax",
            "from langchain_core.prompts import ChatPromptTemplate",
            "from langchain_core.output_parsers import StrOutputParser",
            "from functools import lru_cache",
            "",
            "@lru_cache(maxsize=1)",
            "def _get_analysis_chain():",
            "    prompt_template = ChatPromptTemplate.from_messages([",
            "        ('system', 'You are an expert career consultant.'),",
            "        ('human', '{analysis_prompt}'),",
            "    ])",
            "    # Sprint 5 \u2014 LCEL pipe operator: prompt | llm | output_parser",
            "    return prompt_template | get_llm() | StrOutputParser()",
            "",
            "def analyze_resume_core(resume_text, job_description='', callbacks=None):",
            "    if not resume_text or len(resume_text.strip()) < 50:",
            "        raise ValueError('Resume text too short (min 50 chars).')",
            "    prompt = build_analysis_prompt(resume_text, job_description)",
            "    chain  = _get_analysis_chain()   # cached LCEL chain",
            "    cfg    = {'callbacks': callbacks} if callbacks else {}",
            "    return chain.invoke({'analysis_prompt': prompt}, config=cfg)",
        ],
    ),
    dict(
        num="Sprint 6", topic="Router Chains",
        subtitle="RunnableBranch  \u00b7  intent classifier  \u2192  specialist chain dispatch",
        file="src/backend/agent/router.py",
        accent=RGBColor(0x16, 0xA0, 0x85),
        code=[
            "# agent/router.py \u2014 Sprint 6: Intent routing with RunnableBranch",
            "from langchain_core.runnables import RunnableBranch, RunnablePassthrough",
            "",
            "@lru_cache(maxsize=1)",
            "def get_router_pipeline():",
            "    llm, parser = get_llm(), StrOutputParser()",
            "",
            "    # One specialist LCEL chain per intent (pipe syntax)",
            "    job_search_chain   = _JOB_SEARCH_PROMPT   | llm | parser",
            "    resume_chain       = _RESUME_PROMPT       | llm | parser",
            "    cover_letter_chain = _COVER_LETTER_PROMPT | llm | parser",
            "    general_chain      = _GENERAL_PROMPT      | llm | parser",
            "",
            "    branch = RunnableBranch(",
            "        (lambda x: 'job_search'   in x.get('intent',''), job_search_chain),",
            "        (lambda x: 'resume'       in x.get('intent',''), resume_chain),",
            "        (lambda x: 'cover_letter' in x.get('intent',''), cover_letter_chain),",
            "        general_chain,   # default branch",
            "    )",
            "    classify = _get_classify_chain()   # LLM intent classifier chain",
            "    return RunnablePassthrough.assign(intent=classify) | branch",
        ],
    ),
    dict(
        num="Sprint 7", topic="Streaming (SSE)",
        subtitle="astream_events()  \u00b7  FastAPI StreamingResponse  \u00b7  token-by-token delivery",
        file="src/backend/fastapi_app.py",
        accent=RGBColor(0xC0, 0x39, 0x2B),
        code=[
            "# fastapi_app.py \u2014 Sprint 7: SSE token streaming endpoint",
            "from fastapi.responses import StreamingResponse",
            "import json, uuid",
            "",
            "@app.post('/api/chat/stream')",
            "async def stream_chat(body: ChatRequest, request: Request):",
            "    session_id = body.session_id or str(uuid.uuid4())",
            "",
            "    async def event_generator():",
            "        # stream_agent() yields tokens via .astream_events()",
            "        async for token in stream_agent(body.message, session_id):",
            "            payload = json.dumps({'token': token, 'session_id': session_id})",
            "            yield f'data: {payload}\\n\\n'",
            "        # End-of-stream sentinel for the client",
            "        yield f'data: {json.dumps({\"done\": True})}\\n\\n'",
            "",
            "    return StreamingResponse(",
            "        event_generator(),",
            "        media_type='text/event-stream',",
            "        headers={",
            "            'Cache-Control': 'no-cache',",
            "            'X-Accel-Buffering': 'no',   # disable proxy buffering",
            "        },",
            "    )",
        ],
    ),
    dict(
        num="Sprint 8", topic="Auth0 + Observability",
        subtitle="RS256 JWT  \u00b7  JWKS validation  \u00b7  Langfuse CallbackHandler",
        file="src/backend/auth/auth0.py  |  observability/langfuse_config.py",
        accent=RGBColor(0x27, 0x66, 0x9C),
        code=[
            "# auth/auth0.py \u2014 Sprint 8: Auth0 RS256 JWT verification",
            "@lru_cache(maxsize=1)",
            "def _get_jwks() -> dict:",
            "    resp = http_requests.get(f'https://{AUTH0_DOMAIN}/.well-known/jwks.json')",
            "    resp.raise_for_status()",
            "    return resp.json()   # cached for the process lifetime",
            "",
            "def verify_token(token: str) -> dict:",
            "    jwks   = _get_jwks()",
            "    header = jwt.get_unverified_header(token)",
            "    key    = next((k for k in jwks['keys'] if k['kid']==header['kid']), None)",
            "    return jwt.decode(token, key, algorithms=['RS256'],",
            "                     audience=AUTH0_AUDIENCE)",
            "",
            "# observability/langfuse_config.py \u2014 Langfuse CallbackHandler injection",
            "def get_langfuse_handler(session_id, user_id=None):",
            "    from langfuse.callback import CallbackHandler",
            "    return CallbackHandler(",
            "        secret_key=LANGFUSE_SK, public_key=LANGFUSE_PK,",
            "        session_id=session_id,  user_id=user_id,",
            "        trace_name='job-placement-agent',",
            "        tags=['job-agent', 'langchain'],",
            "    )",
        ],
    ),
    dict(
        num="Sprint 9", topic="gRPC Transport",
        subtitle="Protobuf service  \u00b7  grpcio servicer  \u00b7  gRPC-Web for browser clients",
        file="src/backend/grpc_server.py  +  proto/job_agent.proto",
        accent=RGBColor(0x7F, 0x8C, 0x8D),
        code=[
            "# grpc_server.py \u2014 Sprint 9: gRPC servicer",
            "# Proto: service JobAgentService {",
            "#   rpc HealthCheck(HealthRequest)  returns (HealthResponse);",
            "#   rpc Chat       (ChatRequest)    returns (ChatResponse); }",
            "",
            "class JobAgentServicer(job_agent_pb2_grpc.JobAgentServiceServicer):",
            "    def HealthCheck(self, request, context):",
            "        return job_agent_pb2.HealthResponse(",
            "            status='healthy', version='1.0.0',",
            "            message='gRPC server is running.',",
            "        )",
            "",
            "    def Chat(self, request, context):  # protected RPC",
            "        user       = _require_auth(context)   # RS256 JWT check",
            "        session_id = request.session_id or str(uuid.uuid4())",
            "        try:",
            "            resp = run_agent(request.message, session_id, user['sub'])",
            "            return job_agent_pb2.ChatResponse(",
            "                response=resp, session_id=session_id)",
            "        except AppError as exc:",
            "            code, detail = _map_app_error(exc)  # domain \u2192 gRPC status",
            "            context.abort(code, detail)",
        ],
    ),
]

# ── Build slides ──────────────────────────────────────────────────────────────
prs = Presentation(PPTX_IN)
W, H = prs.slide_width, prs.slide_height

# Insert in reverse so slide 6 stays at index 5
for i in range(7, -1, -1):
    sp = SPRINTS[i]
    accent = sp["accent"]
    slide = insert_blank_slide(prs, after_index=5)

    add_rect(slide, 0, 0, W, H, C_BG)                        # background
    add_rect(slide, 0, 0, Inches(0.15), H, accent)            # left bar

    # Badge
    add_rect(slide, Inches(0.25), Inches(0.18), Inches(1.35), Inches(0.38), accent)
    add_tb(slide, Inches(0.25), Inches(0.19), Inches(1.35), Inches(0.36),
           sp["num"], Pt(14), C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Topic title
    add_tb(slide, Inches(1.75), Inches(0.16), Inches(10.3), Inches(0.45),
           sp["topic"], Pt(24), C_WHITE, bold=True)

    # Subtitle
    add_tb(slide, Inches(1.75), Inches(0.60), Inches(10.3), Inches(0.30),
           sp["subtitle"], Pt(11), C_HEADER, italic=True)

    # Separator
    add_rect(slide, Inches(0.25), Inches(0.95), W - Inches(0.50), Inches(0.025), accent)

    # File path
    add_tb(slide, Inches(0.25), Inches(1.00), Inches(11.5), Inches(0.25),
           sp["file"], Pt(9), RGBColor(0x80, 0xA0, 0xC0), italic=True, font="Courier New")

    # Code block
    code_y = Inches(1.30)
    code_h = H - code_y - Inches(0.18)
    add_code(slide, Inches(0.25), code_y, W - Inches(0.50), code_h, sp["code"])

prs.save(PPTX)
prs2 = Presentation(PPTX)
print(f"Saved. Total slides: {len(prs2.slides)}")
for i, sl in enumerate(prs2.slides):
    texts = [s.text.strip() for s in sl.shapes
             if hasattr(s, "text") and s.text.strip()]
    first = texts[0] if texts else "(empty)"
    print(f"  Slide {i+1:2d}: {first[:60]}")
