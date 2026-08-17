"""
update_ppt_v3.py — Add Guest Mode + Performance Optimisation slides,
                    update Deployment slide with live Vercel URLs.

Run with the system Python (where python-pptx is installed):
    python docs/update_ppt_v3.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy, os

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
SRC    = os.path.join(BASE, "Job_Placement_Agent_Demo_v2.pptx")
DST    = os.path.join(BASE, "Job_Placement_Agent_Demo_v3.pptx")

# ── Colour palette (matches generate_ppt.py) ──────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x1E, 0x88, 0xE5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
AMBER   = RGBColor(0xF5, 0x7F, 0x17)
PURPLE  = RGBColor(0x6A, 0x1B, 0x9A)
TEAL    = RGBColor(0x00, 0x69, 0x6B)
GREY    = RGBColor(0x44, 0x44, 0x44)
LGREY   = RGBColor(0xF5, 0xF5, 0xF5)
MUTED   = RGBColor(0x78, 0x90, 0x9C)
RED     = RGBColor(0xC6, 0x28, 0x28)

# ── Helpers ───────────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill, line=None, alpha=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    MSO_RECT = 1
    shp = slide.shapes.add_shape(MSO_RECT, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


def _txt(slide, text, x, y, w, h,
         size=11, bold=False, italic=False,
         color=GREY, align=PP_ALIGN.LEFT,
         font="Segoe UI", wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb


def _add_para(txb, text, size=10, bold=False, italic=False,
              color=GREY, align=PP_ALIGN.LEFT, font="Segoe UI", space_before=0):
    tf = txb.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return p


def _slide_chrome(slide, prs, title_text, accent_color=ACCENT):
    """Standard chrome: white bg, top accent bar, title."""
    W, H = prs.slide_width, prs.slide_height
    # White background
    _rect(slide, 0, 0, W, H, WHITE)
    # Top accent bar
    _rect(slide, 0, 0, W, Inches(0.08), accent_color)
    # Title
    _txt(slide, title_text,
         Inches(0.45), Inches(0.14), W - Inches(0.9), Inches(0.55),
         size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)


def _card(slide, x, y, w, h, header, header_color, lines, icon=""):
    """Rounded card with coloured header strip and bullet lines."""
    _rect(slide, x, y, w, h, LGREY, line=RGBColor(0xDD, 0xDD, 0xDD))
    _rect(slide, x, y, w, Inches(0.32), header_color)
    # Header text
    label = f"{icon}  {header}" if icon else header
    _txt(slide, label, x + Inches(0.1), y + Inches(0.04),
         w - Inches(0.2), Inches(0.26),
         size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Body lines
    body_x = x + Inches(0.12)
    body_y = y + Inches(0.38)
    body_w = w - Inches(0.24)
    body_h = h - Inches(0.44)
    if lines:
        txb = _txt(slide, lines[0], body_x, body_y, body_w, body_h,
                   size=9.5, color=GREY)
        for line in lines[1:]:
            _add_para(txb, line, size=9.5, color=GREY)


def _insert_slide_at(prs, slide, idx):
    """Move the last added slide to position idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    entries = list(xml_slides)
    # The newly added slide is the last entry
    entry = entries[-1]
    xml_slides.remove(entry)
    xml_slides.insert(idx, entry)


# ── Slide A: Guest Mode & Optional Auth0 ──────────────────────────────────────

def slide_guest_mode(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    W, H   = prs.slide_width, prs.slide_height

    _slide_chrome(slide, prs, "Optional Auth0 & Guest Mode", TEAL)

    # Sub-heading
    _txt(slide,
         "Users can now access the app without signing in — no Auth0 account required",
         Inches(0.45), Inches(0.72), W - Inches(0.9), Inches(0.32),
         size=11, italic=True, color=MUTED)

    # ── Left column — How it works ─────────────────────────────────────────
    col_w = Inches(3.6)
    col_h = Inches(4.4)
    cx    = Inches(0.4)
    cy    = Inches(1.12)

    _rect(slide, cx, cy, col_w, col_h, RGBColor(0xE0, 0xF2, 0xF1),
          line=RGBColor(0x80, 0xCB, 0xC4))
    _rect(slide, cx, cy, col_w, Inches(0.34), TEAL)
    _txt(slide, "⚙  How It Works", cx + Inches(0.12), cy + Inches(0.05),
         col_w - Inches(0.2), Inches(0.26), size=10, bold=True, color=WHITE)

    steps = [
        ("1", "Login screen shows two options:", NAVY),
        ("",  "• \"Sign in to continue\" (Auth0)", GREY),
        ("",  "• \"Continue as Guest\"", TEAL),
        ("2", "Guest sets isGuest = true and", NAVY),
        ("",  "calls enterApp(true) — skipping", GREY),
        ("",  "all Auth0 token fetches.", GREY),
        ("3", "API calls route to /api/chat/stream", NAVY),
        ("",  "— a public endpoint requiring no", GREY),
        ("",  "Bearer token.", GREY),
        ("4", "Sidebar shows \"Guest\" badge +", NAVY),
        ("",  "sign-in icon to upgrade anytime.", GREY),
    ]

    txb = _txt(slide, "", cx + Inches(0.15), cy + Inches(0.42),
               col_w - Inches(0.3), col_h - Inches(0.5),
               size=9.5, color=GREY)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for num, text, clr in steps:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3 if num else 0)
        run = p.add_run()
        run.text = ("  " if not num else "") + text
        run.font.size = Pt(9.5 if not num else 10)
        run.font.bold = bool(num)
        run.font.color.rgb = clr
        run.font.name = "Segoe UI"

    # ── Middle column — Guest vs Auth0 ────────────────────────────────────
    mx  = cx + col_w + Inches(0.25)
    my  = cy
    mw  = Inches(2.65)
    mh  = col_h

    _rect(slide, mx, my, mw, mh, RGBColor(0xF3, 0xE5, 0xF5),
          line=RGBColor(0xCE, 0x93, 0xD8))
    _rect(slide, mx, my, mw, Inches(0.34), PURPLE)
    _txt(slide, "🔐  Comparison", mx + Inches(0.1), my + Inches(0.05),
         mw - Inches(0.2), Inches(0.26), size=10, bold=True, color=WHITE)

    rows = [
        ("Feature",      "Guest",   "Auth0"),
        ("Login",        "None",    "Required"),
        ("Session",      "Shared",  "Per-user"),
        ("Langfuse",     "Partial", "Full trace"),
        ("Token",        "None",    "JWT Bearer"),
        ("Upgrade",      "Always",  "N/A"),
        ("API endpoint", "/stream", "/chat"),
    ]

    ry = my + Inches(0.42)
    for i, (feat, g, a) in enumerate(rows):
        bg = RGBColor(0xEC, 0xE0, 0xF4) if i == 0 else (
             RGBColor(0xF8, 0xF2, 0xFC) if i % 2 == 0 else WHITE)
        _rect(slide, mx + Inches(0.06), ry, mw - Inches(0.12), Inches(0.36), bg)
        row_bold = (i == 0)
        _txt(slide, feat, mx + Inches(0.1), ry + Inches(0.06),
             Inches(1.0), Inches(0.26), size=8.5, bold=row_bold, color=NAVY)
        _txt(slide, g, mx + Inches(1.12), ry + Inches(0.06),
             Inches(0.7), Inches(0.26), size=8.5, bold=row_bold,
             color=TEAL if (not row_bold and g not in ("Guest","None","Partial","N/A","Always","/stream")) else GREY,
             align=PP_ALIGN.CENTER)
        _txt(slide, a, mx + Inches(1.85), ry + Inches(0.06),
             Inches(0.7), Inches(0.26), size=8.5, bold=row_bold,
             color=ACCENT, align=PP_ALIGN.CENTER)
        ry += Inches(0.37)

    # ── Right column — Code snippet ────────────────────────────────────────
    rx  = mx + mw + Inches(0.25)
    ry2 = cy
    rw  = W - rx - Inches(0.35)
    rh  = col_h

    _rect(slide, rx, ry2, rw, rh, RGBColor(0x1E, 0x1E, 0x1E))
    _txt(slide, "index.html — guest flow", rx + Inches(0.1), ry2 + Inches(0.06),
         rw - Inches(0.2), Inches(0.22), size=8, italic=True,
         color=RGBColor(0x85, 0x99, 0xAA), font="Courier New")

    code_lines = [
        ("// State",                               RGBColor(0x6A, 0x99, 0x55)),
        ("let isGuest = false;",                   RGBColor(0xD4, 0xD4, 0xD4)),
        ("",                                       RGBColor(0xD4, 0xD4, 0xD4)),
        ("// Button handler",                      RGBColor(0x6A, 0x99, 0x55)),
        ("async function continueAsGuest() {",     RGBColor(0x56, 0x9C, 0xD6)),
        ("  isGuest = true;",                      RGBColor(0xD4, 0xD4, 0xD4)),
        ("  hide('login-screen');",                RGBColor(0xD4, 0xD4, 0xD4)),
        ("  await enterApp(true);",                RGBColor(0xD4, 0xD4, 0xD4)),
        ("}",                                      RGBColor(0x56, 0x9C, 0xD6)),
        ("",                                       RGBColor(0xD4, 0xD4, 0xD4)),
        ("// Route to public endpoint",            RGBColor(0x6A, 0x99, 0x55)),
        ("const ep = isGuest",                     RGBColor(0xD4, 0xD4, 0xD4)),
        ("  ? '/api/chat/stream'",                 RGBColor(0xCE, 0x91, 0x78)),
        ("  : '/api/chat/stream';",                RGBColor(0xCE, 0x91, 0x78)),
        ("",                                       RGBColor(0xD4, 0xD4, 0xD4)),
        ("// Logout handles guest",                RGBColor(0x6A, 0x99, 0x55)),
        ("async function logout() {",              RGBColor(0x56, 0x9C, 0xD6)),
        ("  if (isGuest) {",                       RGBColor(0xD4, 0xD4, 0xD4)),
        ("    isGuest = false;",                   RGBColor(0xD4, 0xD4, 0xD4)),
        ("    show('login-screen');",              RGBColor(0xD4, 0xD4, 0xD4)),
        ("    return;",                            RGBColor(0xD4, 0xD4, 0xD4)),
        ("  }",                                    RGBColor(0xD4, 0xD4, 0xD4)),
        ("  await auth0Client.logout(...);",       RGBColor(0xD4, 0xD4, 0xD4)),
        ("}",                                      RGBColor(0x56, 0x9C, 0xD6)),
    ]

    cy2 = ry2 + Inches(0.34)
    LINE_H = Inches(0.175)
    for line_text, clr in code_lines:
        if line_text:
            _txt(slide, line_text, rx + Inches(0.1), cy2,
                 rw - Inches(0.2), LINE_H + Inches(0.02),
                 size=8.5, color=clr, font="Courier New", wrap=False)
        cy2 += LINE_H

    # ── Bottom tag ────────────────────────────────────────────────────────
    _txt(slide, "No breaking changes · Auth0 flow unchanged · Public endpoints pre-existed in backend",
         Inches(0.45), H - Inches(0.42), W - Inches(0.9), Inches(0.3),
         size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    return slide


# ── Slide B: Latency & Performance Optimisations ──────────────────────────────

def slide_performance(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    W, H   = prs.slide_width, prs.slide_height

    _slide_chrome(slide, prs, "Latency & Performance Optimisations", AMBER)

    _txt(slide,
         "5 targeted fixes reduced agent response latency and eliminated cold-start blocking",
         Inches(0.45), Inches(0.72), W - Inches(0.9), Inches(0.32),
         size=11, italic=True, color=MUTED)

    # ── 5 fix cards in 2 rows ─────────────────────────────────────────────
    cards = [
        (GREEN,  "1. Async Endpoints",
         ["All sync def routes → async def",
          "Event loop never frozen during LLM/HTTP calls",
          "FastAPI can handle concurrent requests properly"]),
        (ACCENT, "2. asyncio.to_thread()",
         ["SerpAPI, resume analysis, cover letter,",
          "and Auth0 token fetch run in thread pool",
          "Blocking I/O off the event loop"]),
        (TEAL,   "3. SSE Token Streaming",
         ["Frontend uses /api/chat/stream (SSE)",
          "Text appears word-by-word via ReadableStream",
          "Falls back to JSON if streaming fails"]),
        (PURPLE, "4. Startup Pre-warming",
         ["on_startup() launches background task",
          "Auth0 JWKS + Agent executor built at boot",
          "First real request hits warm singletons"]),
        (AMBER,  "5. Langfuse Cache",
         ["_langfuse_setup() decorated with @lru_cache",
          "Env-var reads + package import happen once",
          "Per-request overhead: only CallbackHandler()"]),
    ]

    card_w = Inches(1.88)
    card_h = Inches(2.3)
    gap    = Inches(0.18)
    start_x = Inches(0.38)
    row1_y  = Inches(1.12)

    for i, (clr, title, lines) in enumerate(cards):
        cx = start_x + i * (card_w + gap)
        cy = row1_y

        _rect(slide, cx, cy, card_w, card_h, LGREY,
              line=RGBColor(0xCC, 0xCC, 0xCC))
        _rect(slide, cx, cy, card_w, Inches(0.36), clr)
        _txt(slide, title,
             cx + Inches(0.1), cy + Inches(0.05),
             card_w - Inches(0.2), Inches(0.28),
             size=10, bold=True, color=WHITE)

        body_y = cy + Inches(0.42)
        txb = _txt(slide, lines[0],
                   cx + Inches(0.1), body_y,
                   card_w - Inches(0.2), card_h - Inches(0.5),
                   size=9, color=GREY, wrap=True)
        for ln in lines[1:]:
            _add_para(txb, ln, size=9, color=GREY, space_before=3)

    # ── Before / After table ──────────────────────────────────────────────
    ty  = row1_y + card_h + Inches(0.22)
    tw  = W - Inches(0.76)
    th  = Inches(1.58)
    tx  = Inches(0.38)

    _rect(slide, tx, ty, tw, th, LGREY, line=RGBColor(0xCC, 0xCC, 0xCC))
    _txt(slide, "Response Latency Profile",
         tx + Inches(0.15), ty + Inches(0.07),
         tw, Inches(0.26), size=11, bold=True, color=NAVY)

    headers = ["Scenario", "Before", "After", "Key Fix"]
    col_ws  = [Inches(2.55), Inches(1.35), Inches(1.35), Inches(4.35)]
    hx = tx + Inches(0.1)
    hy = ty + Inches(0.36)
    for h_text, cw in zip(headers, col_ws):
        _rect(slide, hx, hy, cw - Inches(0.04), Inches(0.28), NAVY)
        _txt(slide, h_text, hx + Inches(0.06), hy + Inches(0.04),
             cw - Inches(0.1), Inches(0.22), size=9, bold=True, color=WHITE)
        hx += cw

    rows_data = [
        ("Chat — subsequent call",    "1–3 s",  "1–3 s",  "Async → no event-loop block"),
        ("Chat — cold start",          "4–8 s",  "1–3 s",  "Pre-warm at startup"),
        ("Time-to-first-token",        "full wait","~100 ms","SSE streaming"),
    ]
    clrs = [WHITE, LGREY]
    for ri, (sc, bef, aft, fix) in enumerate(rows_data):
        ry = hy + Inches(0.30) + ri * Inches(0.28)
        vals = [sc, bef, aft, fix]
        vx = tx + Inches(0.1)
        for vi, (vt, cw) in enumerate(zip(vals, col_ws)):
            _rect(slide, vx, ry, cw - Inches(0.04), Inches(0.26), clrs[ri % 2])
            txt_clr = RED if (vi == 1 and ri < 2) else (GREEN if vi == 2 else GREY)
            _txt(slide, vt, vx + Inches(0.06), ry + Inches(0.04),
                 cw - Inches(0.1), Inches(0.2), size=8.5, color=txt_clr,
                 bold=(vi in (1, 2)))
            vx += cw

    # ── Footer ────────────────────────────────────────────────────────────
    _txt(slide,
         "maxDuration: 60 added to vercel.json  ·  BOM-stripping middleware for Windows env-var safety",
         Inches(0.45), H - Inches(0.42), W - Inches(0.9), Inches(0.3),
         size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    return slide


# ── Update Deployment slide (slide 20) ────────────────────────────────────────

def update_deployment_slide(prs):
    slide = prs.slides[19]   # 0-based index → slide 20

    W, H = prs.slide_width, prs.slide_height

    # Add live URL box at the bottom of the existing slide
    url_y = H - Inches(0.72)
    url_h = Inches(0.58)
    _rect(slide, Inches(0.38), url_y, W - Inches(0.76), url_h,
          RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)

    _txt(slide, "🚀  Live Deployment",
         Inches(0.55), url_y + Inches(0.06),
         Inches(1.5), Inches(0.28), size=10, bold=True, color=GREEN)

    urls = (
        "Frontend: https://job-agent-frontend-seven.vercel.app    "
        "│    Backend API: https://job-agent-backend-nine.vercel.app    "
        "│    Docs: https://job-agent-backend-nine.vercel.app/docs"
    )
    _txt(slide, urls,
         Inches(1.9), url_y + Inches(0.08),
         W - Inches(2.3), Inches(0.38),
         size=8.5, color=NAVY, font="Courier New")


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(SRC)

    print(f"Opened: {SRC}  ({len(prs.slides)} slides)")

    # Add two new slides (appended, then repositioned)
    slide_a = slide_guest_mode(prs)
    slide_b = slide_performance(prs)

    # Insert Guest Mode after slide 5 (Key Features, index 4)
    _insert_slide_at(prs, slide_a, 5)   # now slide 6

    # Insert Performance after new slide 6, so index 6 → slide 7
    _insert_slide_at(prs, slide_b, 6)

    # Update deployment slide (was slide 20, now slide 22 after 2 inserts)
    update_deployment_slide(prs)

    prs.save(DST)
    print(f"Saved:  {DST}  ({len(prs.slides)} slides)")
    print()
    print("New slide order:")
    for i, sl in enumerate(prs.slides, 1):
        t = ""
        for sh in sl.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()[:55]
                if t:
                    break
        print(f"  {i:2d}.  {t}")


if __name__ == "__main__":
    main()
