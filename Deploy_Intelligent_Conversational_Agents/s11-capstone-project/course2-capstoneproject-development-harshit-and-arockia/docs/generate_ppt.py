"""
Generates the Job Placement Agent project demo presentation.

Run from the project root:
    python docs/generate_ppt.py

Output: docs/Job_Placement_Agent_Demo.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand colours ─────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT_BLUE = RGBColor(0x1E, 0x88, 0xE5)   # material blue 600
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)   # blue 50
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x37, 0x47, 0x4F)
LIGHT_GREY  = RGBColor(0xEC, 0xEF, 0xF1)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
AMBER       = RGBColor(0xF5, 0x7F, 0x17)

OUTPUT = os.path.join(os.path.dirname(__file__), "Job_Placement_Agent_Demo.pptx")

# ── Helpers ───────────────────────────────────────────────────────────────────

def _add_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _add_bg(slide, DARK_BLUE)

    # Top accent bar
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    # Gradient block for title area
    _add_rect(slide, 0, 1.0, 10, 3.8, RGBColor(0x10, 0x22, 0x35))

    # Badge
    badge = _add_rect(slide, 0.5, 1.1, 2.2, 0.45, ACCENT_BLUE)

    _add_text_box(slide, "SPRINT 11 CAPSTONE", 0.52, 1.13, 2.1, 0.4,
                  font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main title
    _add_text_box(slide, "Job Placement Agent", 0.5, 1.7, 9.0, 1.3,
                  font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    _add_text_box(slide,
                  "An AI-Powered Career Assistant — Job Search · Resume Analysis\n"
                  "Skill Gap Detection · Cover Letter Generation",
                  0.5, 3.1, 9.0, 1.0,
                  font_size=18, bold=False,
                  color=RGBColor(0xB0, 0xBE, 0xC5), align=PP_ALIGN.LEFT)

    # Tech tags
    tags = ["Gemini LLM", "LangChain", "FastAPI", "gRPC", "Auth0", "Langfuse", "Vercel"]
    x = 0.5
    for tag in tags:
        _add_rect(slide, x, 4.3, 1.25, 0.32, DARK_GREY)
        _add_text_box(slide, tag, x + 0.05, 4.32, 1.15, 0.28,
                      font_size=9, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        x += 1.32

    # Bottom divider
    _add_rect(slide, 0, 4.85, 10, 0.04, ACCENT_BLUE)

    # Footer
    _add_text_box(slide, "NIIT — Build and Deploy Intelligent Conversational Agents",
                  0.5, 5.0, 9.0, 0.4,
                  font_size=10, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.LEFT)

    return slide


def _agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)
    _add_rect(slide, 0, 0.08, 0.08, 7.42, ACCENT_BLUE)

    _add_text_box(slide, "Agenda", 0.3, 0.2, 9.5, 0.7,
                  font_size=32, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

    items = [
        ("01", "Problem Statement", "Why this agent matters"),
        ("02", "Architecture Overview", "System design at a glance"),
        ("03", "Key Features", "What the agent can do"),
        ("04", "Sprint Coverage", "Sprints 2-9 concepts applied"),
        ("05", "Live Demo Flow", "End-to-end user journey"),
        ("06", "Observability & Auth", "Langfuse tracing · Auth0 JWT"),
        ("07", "Testing & Quality", "99% test coverage · 376 tests"),
        ("08", "Deployment", "Vercel serverless · gRPC + FastAPI"),
    ]

    col_w = 4.6
    for i, (num, title, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.3 + col * col_w
        y = 1.1 + row * 1.4

        _add_rect(slide, x, y, col_w - 0.2, 1.2, LIGHT_BLUE)
        _add_rect(slide, x, y, 0.55, 1.2, ACCENT_BLUE)
        _add_text_box(slide, num, x + 0.04, y + 0.25, 0.47, 0.6,
                      font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, title, x + 0.65, y + 0.1, col_w - 0.85, 0.45,
                      font_size=14, bold=True, color=DARK_BLUE)
        _add_text_box(slide, sub, x + 0.65, y + 0.6, col_w - 0.85, 0.45,
                      font_size=10, color=DARK_GREY)

    return slide


def _problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Problem Statement", 0.4, 0.2, 9.5, 0.7,
                  font_size=30, bold=True, color=DARK_BLUE)

    _add_text_box(slide,
                  "Job seekers waste hours on fragmented, manual tasks:",
                  0.4, 1.0, 9.2, 0.5, font_size=16, color=DARK_GREY)

    problems = [
        ("Searching dozens of sites", "No single place to find relevant roles by title + location"),
        ("Unclear resume gaps",       "Hard to know which skills are missing without expert feedback"),
        ("Generic cover letters",     "Tailoring letters for each role is time-consuming"),
        ("No guided workflow",        "Users jump between tools with no connected, intelligent flow"),
    ]

    for i, (title, desc) in enumerate(problems):
        y = 1.65 + i * 1.1
        _add_rect(slide, 0.4, y, 0.06, 0.7, ACCENT_BLUE)
        _add_text_box(slide, title, 0.6, y,      8.8, 0.38, font_size=14, bold=True,  color=DARK_BLUE)
        _add_text_box(slide, desc,  0.6, y + 0.38, 8.8, 0.38, font_size=11, color=DARK_GREY)

    _add_rect(slide, 0.4, 6.1, 9.2, 0.5, ACCENT_BLUE)
    _add_text_box(slide,
                  "Solution: One conversational AI agent that guides users from search → analysis → application",
                  0.55, 6.15, 9.0, 0.4, font_size=13, bold=True, color=WHITE)

    return slide


def _architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Architecture Overview", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    # Layer 1 — Frontend
    _add_rect(slide, 0.3, 1.0, 9.4, 0.75, LIGHT_BLUE)
    _add_rect(slide, 0.3, 1.0, 0.12, 0.75, ACCENT_BLUE)
    _add_text_box(slide, "FRONTEND  (Streamlit / index.html)", 0.55, 1.05, 5.0, 0.3,
                  font_size=11, bold=True, color=DARK_BLUE)
    _add_text_box(slide, "Chat UI · Resume Paste · Quick Actions · Session Management",
                  0.55, 1.38, 8.5, 0.3, font_size=9, color=DARK_GREY)

    # Arrow
    _add_text_box(slide, "HTTP REST / gRPC-Web", 4.2, 1.82, 3.0, 0.3,
                  font_size=8, color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Layer 2 — FastAPI Backend
    _add_rect(slide, 0.3, 2.1, 9.4, 0.85, LIGHT_BLUE)
    _add_rect(slide, 0.3, 2.1, 0.12, 0.85, ACCENT_BLUE)
    _add_text_box(slide, "BACKEND  (FastAPI · gRPC Server · Vercel Serverless)",
                  0.55, 2.14, 6.0, 0.3, font_size=11, bold=True, color=DARK_BLUE)
    _add_text_box(slide,
                  "/api/health · /api/auth/token · /api/chat · /api/chat/stream · "
                  "/api/chat/route · /api/jobs/search · /api/resume/analyze · /api/cover-letter/gen",
                  0.55, 2.48, 8.5, 0.35, font_size=8, color=DARK_GREY)

    # Auth0 box inline
    _add_rect(slide, 7.5, 2.14, 2.1, 0.65, RGBColor(0xFF, 0xF3, 0xE0))
    _add_text_box(slide, "Auth0 JWT\nRS256 Validation", 7.55, 2.18, 2.0, 0.58,
                  font_size=8, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # Layer 3 — LangChain Agent
    _add_rect(slide, 0.3, 3.1, 9.4, 0.85, LIGHT_BLUE)
    _add_rect(slide, 0.3, 3.1, 0.12, 0.85, GREEN)
    _add_text_box(slide, "LANGCHAIN AGENT  (job_agent.py)",
                  0.55, 3.14, 6.0, 0.3, font_size=11, bold=True, color=DARK_BLUE)
    _add_text_box(slide,
                  "create_tool_calling_agent · AgentExecutor · RunnableWithMessageHistory · "
                  "LCEL Router · Async .ainvoke() · SSE astream_events()",
                  0.55, 3.48, 8.5, 0.35, font_size=8, color=DARK_GREY)

    # Layer 4 — Tools
    tools = [
        ("search_jobs", "SerpAPI\nGoogle Jobs", RGBColor(0xE8, 0xF5, 0xE9)),
        ("analyze_resume", "LCEL Chain\nStrOutputParser", RGBColor(0xE8, 0xF5, 0xE9)),
        ("generate_cover_letter", "LCEL Chain\nStrOutputParser", RGBColor(0xE8, 0xF5, 0xE9)),
    ]
    for i, (name, detail, bg) in enumerate(tools):
        x = 0.3 + i * 3.15
        _add_rect(slide, x, 4.1, 3.0, 0.85, bg)
        _add_text_box(slide, name,   x + 0.1, 4.12, 2.8, 0.32, font_size=10, bold=True, color=GREEN)
        _add_text_box(slide, detail, x + 0.1, 4.45, 2.8, 0.42, font_size=8, color=DARK_GREY)

    # Layer 5 — External services
    services = [
        ("Gemini LLM", "gemini-2.0-flash\nOpenAI-compat API", ACCENT_BLUE),
        ("SerpAPI",    "Real-time Google\nJobs results",       ACCENT_BLUE),
        ("Langfuse",   "LLM Observability\nCallbackHandler",   ACCENT_BLUE),
    ]
    for i, (name, detail, color) in enumerate(services):
        x = 0.3 + i * 3.15
        _add_rect(slide, x, 5.1, 3.0, 0.75, DARK_BLUE)
        _add_text_box(slide, name,   x + 0.1, 5.12, 2.8, 0.3, font_size=10, bold=True, color=WHITE)
        _add_text_box(slide, detail, x + 0.1, 5.43, 2.8, 0.35, font_size=8, color=RGBColor(0xB0, 0xBE, 0xC5))

    return slide


def _features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Key Features", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    features = [
        ("Real-Time Job Search",       "SerpAPI Google Jobs — 5-10 live listings by role + city",          ACCENT_BLUE),
        ("Resume Analysis",            "LLM-powered skill extraction, experience scoring, quality rating",  ACCENT_BLUE),
        ("Skill Gap Detection",        "Side-by-side resume vs. JD comparison with match percentage",       ACCENT_BLUE),
        ("Cover Letter Generator",     "Tailored, role-specific letters generated in seconds",              ACCENT_BLUE),
        ("Conversational Memory",      "RunnableWithMessageHistory — context preserved across turns",        GREEN),
        ("SSE Token Streaming",        "Token-by-token response via FastAPI StreamingResponse",              GREEN),
        ("LCEL Intent Router",         "RunnableBranch classifies intent → routes to specialist chain",     GREEN),
        ("Async Agent (.ainvoke)",     "Non-blocking LLM calls keep the FastAPI event loop free",           GREEN),
        ("Auth0 JWT Auth",             "RS256 signature validation via JWKS — all protected endpoints",     AMBER),
        ("Langfuse Observability",     "Full trace: LLM calls, tools, latency, errors, session grouping",   AMBER),
        ("99% Test Coverage",          "376 tests — unit, integration, FastAPI TestClient, gRPC servicer",  DARK_GREY),
        ("Dual Transport",             "FastAPI (REST + SSE) + gRPC + gRPC-Web on a single servicer",       DARK_GREY),
    ]

    col_w = 4.6
    for i, (title, desc, color) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.3 + col * col_w
        y = 1.0 + row * 1.05

        _add_rect(slide, x, y, 0.06, 0.75, color)
        _add_text_box(slide, title, x + 0.15, y,        col_w - 0.25, 0.35, font_size=11, bold=True, color=DARK_BLUE)
        _add_text_box(slide, desc,  x + 0.15, y + 0.36, col_w - 0.25, 0.38, font_size=9,  color=DARK_GREY)

    return slide


def _sprints_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Sprint Coverage (2 – 9)", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    sprints = [
        ("Sprint 2", "LLM Integration",        "ChatOpenAI via Gemini OpenAI-compat API · lru_cache · Async .ainvoke()"),
        ("Sprint 3", "LangChain Tools",         "search_jobs · analyze_resume · generate_cover_letter — @tool decorator"),
        ("Sprint 4", "Agent + Memory",          "create_tool_calling_agent · AgentExecutor · RunnableWithMessageHistory"),
        ("Sprint 5", "LCEL Chains",             "Pipe syntax  prompt | llm | StrOutputParser()  in every tool"),
        ("Sprint 6", "Router Chains",           "RunnableBranch · intent classifier → specialist chain dispatch"),
        ("Sprint 7", "Streaming (SSE)",         "astream_events() · FastAPI StreamingResponse · token-by-token delivery"),
        ("Sprint 8", "Auth0 + Observability",   "RS256 JWT · JWKS validation · Langfuse CallbackHandler"),
        ("Sprint 9", "gRPC Transport",          "Protobuf service · grpcio servicer · gRPC-Web for browser clients"),
    ]

    for i, (sprint, title, detail) in enumerate(sprints):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 4.85
        y = 1.0 + row * 1.35

        _add_rect(slide, x, y, 4.6, 1.15, LIGHT_BLUE)
        _add_rect(slide, x, y, 0.95, 1.15, ACCENT_BLUE)
        _add_text_box(slide, sprint, x + 0.08, y + 0.35, 0.78, 0.45,
                      font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, title,  x + 1.05, y + 0.05, 3.45, 0.38,
                      font_size=12, bold=True, color=DARK_BLUE)
        _add_text_box(slide, detail, x + 1.05, y + 0.48, 3.45, 0.58,
                      font_size=8, color=DARK_GREY)

    return slide


def _demo_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Live Demo — End-to-End User Journey", 0.4, 0.2, 9.5, 0.6,
                  font_size=28, bold=True, color=DARK_BLUE)

    steps = [
        ("1", "Session Start",
         "User opens chat UI → POST /api/session → server returns UUID session_id"),
        ("2", "Job Discovery",
         '"Find Python developer jobs in Austin" → Agent calls search_jobs() via SerpAPI → 8 live listings returned'),
        ("3", "Resume Analysis",
         'User pastes resume → Agent calls analyze_resume() via LCEL chain → Skills, gaps, score returned'),
        ("4", "Intent Routing",
         "POST /api/chat/route → LCEL RunnableBranch classifies 'cover letter' intent → specialist chain"),
        ("5", "Cover Letter",
         "Agent calls generate_cover_letter() → tailored 350-word letter in seconds"),
        ("6", "Streaming Reply",
         "POST /api/chat/stream → SSE StreamingResponse → tokens appear progressively in UI"),
    ]

    for i, (num, title, detail) in enumerate(steps):
        row = i // 2
        col = i % 2
        x = 0.3 + col * 4.85
        y = 1.0 + row * 1.75

        _add_rect(slide, x, y, 4.6, 1.55, LIGHT_GREY)
        _add_rect(slide, x, y, 0.5, 0.5, ACCENT_BLUE)
        _add_text_box(slide, num, x + 0.1, y + 0.08, 0.32, 0.35,
                      font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, title,  x + 0.6, y + 0.08, 3.9, 0.38,
                      font_size=12, bold=True, color=DARK_BLUE)
        _add_text_box(slide, detail, x + 0.15, y + 0.58, 4.35, 0.85,
                      font_size=9, color=DARK_GREY)

    return slide


def _api_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "API Endpoints", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    # Headers
    _add_rect(slide, 0.3, 1.0, 9.4, 0.35, DARK_BLUE)
    for pos, label in [(0.35, "Method"), (1.1, "Endpoint"), (4.6, "Auth"), (5.4, "Description")]:
        _add_text_box(slide, label, pos, 1.05, 1.2, 0.25,
                      font_size=9, bold=True, color=WHITE)

    endpoints = [
        ("GET",    "/api/health",             "Public",  "Health check"),
        ("POST",   "/api/auth/token",         "Public",  "Get Auth0 access token"),
        ("POST",   "/api/chat/public",        "Public",  "Chat (no auth — demo/Streamlit)"),
        ("POST",   "/api/session",            "Public",  "Create a new session UUID"),
        ("POST",   "/api/chat",               "JWT",     "Chat with the agent (authenticated)"),
        ("POST",   "/api/chat/async",         "JWT",     "Async chat via .ainvoke()"),
        ("POST",   "/api/chat/stream",        "Public",  "SSE streaming — token-by-token"),
        ("POST",   "/api/chat/route",         "Public",  "LCEL intent router dispatch"),
        ("POST",   "/api/jobs/search",        "JWT",     "Direct SerpAPI job search"),
        ("POST",   "/api/resume/analyze",     "JWT",     "Direct resume analysis with Langfuse"),
        ("POST",   "/api/cover-letter/gen",   "JWT",     "Direct cover letter generation"),
        ("DELETE", "/api/chat/{session_id}",  "JWT",     "Clear session conversation memory"),
        ("GET",    "/api/sessions",           "JWT",     "List all active session IDs"),
    ]

    method_colors = {
        "GET":    RGBColor(0x1B, 0x5E, 0x20),
        "POST":   RGBColor(0x0D, 0x47, 0xA1),
        "DELETE": RGBColor(0xB7, 0x1C, 0x1C),
    }

    for i, (method, path, auth, desc) in enumerate(endpoints):
        y = 1.4 + i * 0.38
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        _add_rect(slide, 0.3, y, 9.4, 0.35, bg)

        _add_rect(slide, 0.32, y + 0.05, 0.65, 0.24, method_colors.get(method, DARK_GREY))
        _add_text_box(slide, method, 0.33, y + 0.06, 0.62, 0.22,
                      font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _add_text_box(slide, path, 1.06, y + 0.05, 3.4, 0.25, font_size=9, color=DARK_BLUE)
        auth_color = AMBER if auth == "JWT" else GREEN
        _add_text_box(slide, auth, 4.58, y + 0.05, 0.7, 0.25, font_size=8,
                      bold=True, color=auth_color)
        _add_text_box(slide, desc, 5.38, y + 0.05, 4.2, 0.25, font_size=8, color=DARK_GREY)

    return slide


def _observability_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Observability & Authentication", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    # Langfuse left panel
    _add_rect(slide, 0.3, 1.0, 4.55, 5.55, LIGHT_BLUE)
    _add_rect(slide, 0.3, 1.0, 4.55, 0.42, ACCENT_BLUE)
    _add_text_box(slide, "Langfuse Observability", 0.45, 1.04, 4.2, 0.35,
                  font_size=13, bold=True, color=WHITE)

    lf_items = [
        "LLM calls — prompts, completions, token counts",
        "Tool invocations — tool name, args, output",
        "Latency — end-to-end and per-step timing",
        "Errors — exceptions with full stack context",
        "Traces grouped by session_id",
        "Named: job-placement-agent",
        "Flush on every request (Vercel-safe)",
        "Enabled only when LANGFUSE_SECRET_KEY is set",
    ]
    for i, item in enumerate(lf_items):
        y = 1.55 + i * 0.6
        _add_rect(slide, 0.42, y + 0.1, 0.08, 0.3, ACCENT_BLUE)
        _add_text_box(slide, item, 0.6, y + 0.05, 4.1, 0.45, font_size=10, color=DARK_GREY)

    # Auth0 right panel
    _add_rect(slide, 5.15, 1.0, 4.55, 5.55, RGBColor(0xFF, 0xF8, 0xE1))
    _add_rect(slide, 5.15, 1.0, 4.55, 0.42, AMBER)
    _add_text_box(slide, "Auth0 JWT Authentication", 5.3, 1.04, 4.2, 0.35,
                  font_size=13, bold=True, color=WHITE)

    auth_items = [
        "RS256 JWT signed by Auth0",
        "JWKS endpoint for public key fetch",
        "Validates: signature, audience, issuer, expiry",
        "lru_cache on JWKS (avoids repeated fetches)",
        "Public endpoints: /api/health, /api/chat/public",
        "Protected: all others require Bearer token",
        "POST /api/auth/token returns access token",
        "python-jose[cryptography] for verification",
    ]
    for i, item in enumerate(auth_items):
        y = 1.55 + i * 0.6
        _add_rect(slide, 5.27, y + 0.1, 0.08, 0.3, AMBER)
        _add_text_box(slide, item, 5.45, y + 0.05, 4.1, 0.45, font_size=10, color=DARK_GREY)

    return slide


def _testing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Testing & Code Quality", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    # Big metric
    _add_rect(slide, 0.3, 1.0, 9.4, 1.2, DARK_BLUE)
    _add_text_box(slide, "99%", 0.6, 1.05, 2.2, 1.1,
                  font_size=60, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, "Total Coverage", 2.9, 1.15, 3.5, 0.5,
                  font_size=20, bold=True, color=WHITE)
    _add_text_box(slide, "376 tests · 0 failures · 1 016 statements",
                  2.9, 1.65, 3.5, 0.4, font_size=12, color=RGBColor(0xB0, 0xBE, 0xC5))
    _add_text_box(slide, "pytest + pytest-cov\npytest-mock", 6.5, 1.15, 3.0, 0.8,
                  font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    test_files = [
        ("test_agent.py",            "52 tests", "run_agent · async · stream · _build_agent_executor · flush swallowing"),
        ("test_llm.py",              "22 tests", "get_llm · invoke_llm · openai isinstance branches · ImportError fallback"),
        ("test_router.py",           "27 tests", "route_query · RunnableBranch · classify chain · intent labels"),
        ("test_tools_*.py",          "31 tests", "job_search · resume_analyzer · cover_letter · GeminiError branches"),
        ("test_fastapi_app.py",      "79 tests", "TestClient · all 13 endpoints · AppError handler · SSE stream"),
        ("test_main.py (gRPC)",      "95 tests", "JobAgentServicer · RPCs · _require_auth · _map_app_error · serve()"),
        ("test_auth0.py",            "38 tests", "fetch_token · verify_token · JWKS cache · error handling"),
        ("test_schemas.py",          "32 tests", "Pydantic models · validators · ApiResponse envelope"),
    ]

    for i, (name, count, desc) in enumerate(test_files):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 4.85
        y = 2.45 + row * 1.05

        _add_rect(slide, x, y, 4.6, 0.9, LIGHT_GREY)
        _add_rect(slide, x + 3.5, y, 1.0, 0.9, ACCENT_BLUE)
        _add_text_box(slide, name,  x + 0.12, y + 0.06, 3.3, 0.35, font_size=10, bold=True, color=DARK_BLUE)
        _add_text_box(slide, desc,  x + 0.12, y + 0.48, 3.3, 0.35, font_size=8, color=DARK_GREY)
        _add_text_box(slide, count, x + 3.52, y + 0.28, 0.95, 0.35,
                      font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return slide


def _deployment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "Deployment Architecture", 0.4, 0.2, 9.5, 0.6,
                  font_size=30, bold=True, color=DARK_BLUE)

    # Vercel box
    _add_rect(slide, 0.3, 1.0, 4.55, 2.8, LIGHT_BLUE)
    _add_rect(slide, 0.3, 1.0, 4.55, 0.42, DARK_BLUE)
    _add_text_box(slide, "Frontend — Vercel (Static)", 0.45, 1.05, 4.2, 0.32,
                  font_size=12, bold=True, color=WHITE)
    items = [
        "index.html + build.js (Vite-bundled)",
        "Environment vars via inject-env.ps1",
        "CORS-enabled for API calls",
        "vercel.json routes all to index.html",
    ]
    for i, item in enumerate(items):
        _add_text_box(slide, f"• {item}", 0.45, 1.55 + i * 0.48, 4.2, 0.42,
                      font_size=10, color=DARK_GREY)

    # Backend Vercel
    _add_rect(slide, 5.15, 1.0, 4.55, 2.8, LIGHT_BLUE)
    _add_rect(slide, 5.15, 1.0, 4.55, 0.42, DARK_BLUE)
    _add_text_box(slide, "Backend — Vercel (Serverless)", 5.3, 1.05, 4.2, 0.32,
                  font_size=12, bold=True, color=WHITE)
    items2 = [
        "FastAPI wrapped by Mangum (ASGI → Lambda)",
        "index.py as Vercel Python handler",
        "vercel.json → src/backend/index.py",
        "ENV vars in Vercel dashboard",
    ]
    for i, item in enumerate(items2):
        _add_text_box(slide, f"• {item}", 5.3, 1.55 + i * 0.48, 4.2, 0.42,
                      font_size=10, color=DARK_GREY)

    # gRPC local
    _add_rect(slide, 0.3, 4.05, 4.55, 2.5, RGBColor(0xE8, 0xF5, 0xE9))
    _add_rect(slide, 0.3, 4.05, 4.55, 0.42, GREEN)
    _add_text_box(slide, "gRPC — Local / Docker", 0.45, 4.09, 4.2, 0.32,
                  font_size=12, bold=True, color=WHITE)
    items3 = [
        "grpc_server.py  — native gRPC :50051",
        "grpc_web_server.py — gRPC-Web :8080",
        "main.py starts both (threaded)",
        "Protobuf generated by generate_proto.py",
    ]
    for i, item in enumerate(items3):
        _add_text_box(slide, f"• {item}", 0.45, 4.58 + i * 0.46, 4.2, 0.4,
                      font_size=10, color=DARK_GREY)

    # CI
    _add_rect(slide, 5.15, 4.05, 4.55, 2.5, RGBColor(0xF3, 0xE5, 0xF5))
    _add_rect(slide, 5.15, 4.05, 4.55, 0.42, RGBColor(0x6A, 0x1B, 0x9A))
    _add_text_box(slide, "Quality Gates", 5.3, 4.09, 4.2, 0.32,
                  font_size=12, bold=True, color=WHITE)
    items4 = [
        "pytest --cov → 99% total coverage",
        "pytest.ini: pythonpath = src/backend",
        "conftest.py: sys.modules stubs (no real LLM)",
        "deploy-vercel.ps1: deploy automation",
    ]
    for i, item in enumerate(items4):
        _add_text_box(slide, f"• {item}", 5.3, 4.58 + i * 0.46, 4.2, 0.4,
                      font_size=10, color=DARK_GREY)

    return slide


def _summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BLUE)
    _add_rect(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    _add_text_box(slide, "What We Built", 0.5, 0.5, 9.0, 0.7,
                  font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    summary = [
        ("Full-Stack AI Agent",  "FastAPI REST + gRPC + gRPC-Web transports on one servicer",          ACCENT_BLUE),
        ("8 Sprint Concepts",    "LCEL · Router · Async · SSE · Auth0 · Langfuse · gRPC · Tools",       ACCENT_BLUE),
        ("Production Ready",     "99% test coverage · typed schemas · domain exceptions · Vercel CI",   GREEN),
        ("Observable & Secure",  "Langfuse traces every LLM call · Auth0 RS256 on all protected routes", AMBER),
    ]

    for i, (title, desc, color) in enumerate(summary):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 4.85
        y = 1.5 + row * 1.85

        _add_rect(slide, x, y, 4.55, 1.65, RGBColor(0x15, 0x2B, 0x3D))
        _add_rect(slide, x, y, 4.55, 0.06, color)
        _add_text_box(slide, title, x + 0.2, y + 0.2, 4.2, 0.5,
                      font_size=18, bold=True, color=color)
        _add_text_box(slide, desc, x + 0.2, y + 0.75, 4.2, 0.75,
                      font_size=11, color=RGBColor(0xB0, 0xBE, 0xC5))

    _add_rect(slide, 0, 5.3, 10, 0.04, ACCENT_BLUE)
    _add_text_box(slide, "Thank you  |  Questions?",
                  0.5, 5.45, 9.0, 0.5,
                  font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, "Job Placement Agent — NIIT Sprint 11 Capstone",
                  0.5, 6.0, 9.0, 0.4,
                  font_size=11, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER)

    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)
    _agenda_slide(prs)
    _problem_slide(prs)
    _architecture_slide(prs)
    _features_slide(prs)
    _sprints_slide(prs)
    _demo_flow_slide(prs)
    _api_slide(prs)
    _observability_slide(prs)
    _testing_slide(prs)
    _deployment_slide(prs)
    _summary_slide(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
